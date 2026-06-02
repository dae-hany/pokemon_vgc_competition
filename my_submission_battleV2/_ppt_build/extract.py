from pptx import Presentation
import io
p = Presentation(r"C:\Users\daeho\pokemon-vgc-engine\my_submission_battleV2\_ppt_build\DaehoV2_design.pptx")
out = io.StringIO()
out.write(f"slides: {len(p.slides)}\n")
for i, s in enumerate(p.slides, 1):
    out.write(f"\n## Slide {i}\n")
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                out.write("  | " + t.replace("\n", " / ") + "\n")
open(r"C:\Users\daeho\pokemon-vgc-engine\my_submission_battleV2\_ppt_build\_extract.txt", "w", encoding="utf-8").write(out.getvalue())
print("written _extract.txt")
